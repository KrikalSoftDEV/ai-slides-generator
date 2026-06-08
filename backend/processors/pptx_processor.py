from pptx import Presentation
from typing import List, Dict, Tuple
from io import BytesIO

from processors.image_processor import extract_text_from_image, update_image_with_text


def process_pptx(content: bytes) -> Tuple[List[Dict], Presentation]:
    """
    Extract text from PPTX file.
    
    Returns:
        Tuple of (text_data, presentation_object)
        text_data: List of dicts containing slide text with slide numbers
        presentation_object: The Presentation object for later modifications
    """
    try:
        prs = Presentation(BytesIO(content))
    except Exception as e:
        raise ValueError(f"Invalid or corrupted PPTX file: {str(e)}")
    
    text_data = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = f"Slide {slide_num}:\n"
        text_blocks = []
        
        for shape_idx, shape in enumerate(slide.shapes):
            if hasattr(shape, "text") and shape.text.strip():
                text_blocks.append({
                    "source": "shape",
                    "shape_index": shape_idx,
                    "text": shape.text.strip(),
                })
            elif getattr(shape, "shape_type", None) == 13:
                ocr_result = extract_text_from_image(shape.image.blob)
                if not ocr_result.get("success"):
                    continue

                for block_idx, block in enumerate(ocr_result.get("text_blocks", [])):
                    text = block.get("text", "").strip()
                    if not text:
                        continue
                    text_blocks.append({
                        "source": "picture",
                        "shape_index": shape_idx,
                        "block_index": block_idx,
                        "text": text,
                        "confidence": block.get("confidence"),
                        "bbox": block.get("bbox"),
                        "image_width": ocr_result.get("image_width"),
                        "image_height": ocr_result.get("image_height"),
                        "content_type": shape.image.content_type,
                    })
        
        if text_blocks:
            full_text = "\n".join(block["text"] for block in text_blocks)
            slide_text += full_text
            
            text_data.append({
                "slide_number": slide_num,
                "text": full_text,
                "text_blocks": text_blocks,
                "full_text_with_label": slide_text
            })
    
    return text_data, prs


def update_pptx_text(prs: Presentation, updated_texts: Dict[int, str | List[Dict]]) -> bytes:
    """
    Update text in PPTX presentation.
    
    Args:
        prs: Presentation object
        updated_texts: Dict mapping slide_number to new text
    
    Returns:
        Bytes of the updated PPTX file
    """
    for slide_num, new_text in updated_texts.items():
        if slide_num < 1 or slide_num > len(prs.slides):
            continue
        
        slide = prs.slides[slide_num - 1]
        if isinstance(new_text, list):
            for block in new_text:
                try:
                    shape_index = int(block.get("shape_index"))
                except (TypeError, ValueError):
                    continue

                if shape_index < 0 or shape_index >= len(slide.shapes):
                    continue

                shape = slide.shapes[shape_index]
                source = block.get("source", "shape")
                if source == "picture" or block.get("bbox"):
                    continue

                if hasattr(shape, "text_frame"):
                    _replace_shape_text(shape, str(block.get("text", "")))

            _update_picture_text(slide, new_text)
            continue

        text_blocks = str(new_text).split("\n")
        block_idx = 0

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() and block_idx < len(text_blocks):
                if hasattr(shape, "text_frame"):
                    _replace_shape_text(shape, text_blocks[block_idx])
                    block_idx += 1
    
    # Save to bytes
    buffer = BytesIO()
    prs.save(buffer)
    output = buffer.getvalue()
    buffer.close()
    
    return output


def _update_picture_text(slide, blocks: List[Dict]) -> None:
    picture_blocks: Dict[int, List[Dict]] = {}
    for block in blocks:
        if block.get("source") != "picture" and not block.get("bbox"):
            continue

        original_text = str(block.get("originalText", block.get("text", "")))
        if str(block.get("text", "")) == original_text:
            continue

        try:
            shape_index = int(block.get("shape_index"))
        except (TypeError, ValueError):
            continue

        picture_blocks.setdefault(shape_index, []).append(block)

    for shape_index, changed_blocks in picture_blocks.items():
        if shape_index < 0 or shape_index >= len(slide.shapes):
            continue

        shape = slide.shapes[shape_index]
        if getattr(shape, "shape_type", None) != 13:
            continue

        r_id = shape._pic.blipFill.blip.rEmbed
        image_part = shape.part.related_part(r_id)
        updated_image, _ = update_image_with_text(
            image_part.blob,
            shape.image.content_type,
            {"blocks": changed_blocks},
        )
        image_part._blob = updated_image


def _replace_shape_text(shape, new_text: str) -> None:
    text_frame = shape.text_frame

    first_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
    first_run = None
    if first_paragraph and first_paragraph.runs:
        first_run = first_paragraph.runs[0]

    text_frame.clear()
    paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()

    if first_run is None:
        paragraph.text = new_text
        return

    run = paragraph.add_run()
    run.text = new_text
    run.font.name = first_run.font.name
    run.font.size = first_run.font.size
    run.font.bold = first_run.font.bold
    run.font.italic = first_run.font.italic
    if first_run.font.color and first_run.font.color.type:
        try:
            run.font.color.rgb = first_run.font.color.rgb
        except AttributeError:
            pass
