import io
import os
from typing import Any, Dict, List, Tuple

import fitz
from PIL import Image, ImageSequence
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from processors.image_processor import (
    _bbox_bounds,
    _block_alignment,
    _expanded_text_region,
    _line_count_from_block,
    _sample_block_colors,
    extract_text_from_image,
    update_image_with_text,
)
from processors.pptx_processor import _replace_shape_text, process_pptx


PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
PDF_RENDER_SCALE = 2.0
DEFAULT_LANDSCAPE_WIDTH_IN = 13.333
DEFAULT_PORTRAIT_HEIGHT_IN = 13.333
EMU_PER_INCH = 914400

IMAGE_CONTENT_TYPES = {
    "image/png",
    "image/jpeg",
    "image/jpg",
    "image/webp",
    "image/gif",
}

PDF_CONTENT_TYPES = {"application/pdf"}

PPTX_CONTENT_TYPES = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def detect_source_kind(filename: str, content_type: str) -> str:
    filename_lower = (filename or "").lower()
    content_type = (content_type or "").lower()

    if filename_lower.endswith(".pdf") or content_type in PDF_CONTENT_TYPES:
        return "pdf"
    if filename_lower.endswith(".ppt"):
        raise ValueError("Legacy .ppt files are not editable here. Save the file as .pptx and upload it again.")
    if filename_lower.endswith(".pptx") or content_type in PPTX_CONTENT_TYPES:
        return "pptx"
    if filename_lower.endswith((".jpg", ".jpeg", ".png", ".webp", ".gif")) or content_type in IMAGE_CONTENT_TYPES:
        return "image"

    raise ValueError("Unsupported file type. Upload JPG, JPEG, PNG, PDF, or PPTX.")


def editable_pptx_filename(filename: str) -> str:
    base = os.path.splitext(os.path.basename(filename or "editable-slides"))[0]
    safe = "".join(char for char in base if char.isalnum() or char in " -_.").strip()
    return f"{(safe or 'editable-slides')[:90]}-editable.pptx"


def extract_file_for_editing(content: bytes, filename: str, content_type: str) -> Dict[str, Any]:
    source_kind = detect_source_kind(filename, content_type)

    if source_kind == "image":
        image_bytes, normalized_type, image_width, image_height = _normalize_image_bytes(content, content_type)
        ocr_result = extract_text_from_image(image_bytes)
        if not ocr_result.get("success"):
            raise ValueError(ocr_result.get("error") or "OCR failed for the uploaded image.")

        return {
            "file_type": "image",
            "filename": filename,
            "content_type": normalized_type,
            "content": {
                "full_text": ocr_result.get("full_text", ""),
                "text_blocks": _prepare_blocks_for_client(
                    ocr_result.get("text_blocks", []),
                    image_width,
                    image_height,
                ),
                "image_width": image_width,
                "image_height": image_height,
            },
            "success": True,
        }

    if source_kind == "pdf":
        pages = _extract_pdf_pages(content)
        return {
            "file_type": "pdf",
            "filename": filename,
            "content": [_page_to_client_slide(page, page_number) for page_number, page in enumerate(pages, start=1)],
            "success": True,
        }

    text_data, _ = process_pptx(content)
    return {
        "file_type": "pptx",
        "filename": filename,
        "content": text_data,
        "success": True,
    }


def convert_file_to_editable_pptx(
    content: bytes,
    filename: str,
    content_type: str,
    updated_content: Any | None = None,
) -> bytes:
    source_kind = detect_source_kind(filename, content_type)

    if source_kind == "image":
        image_bytes, normalized_type, image_width, image_height = _normalize_image_bytes(content, content_type)
        blocks = _image_blocks_from_updated(updated_content)
        if blocks is None:
            ocr_result = extract_text_from_image(image_bytes)
            if not ocr_result.get("success"):
                raise ValueError(ocr_result.get("error") or "OCR failed for the uploaded image.")
            blocks = ocr_result.get("text_blocks", [])

        page = {
            "label": filename or "Uploaded image",
            "image_bytes": image_bytes,
            "content_type": normalized_type,
            "image_width": image_width,
            "image_height": image_height,
            "blocks": _prepare_blocks_for_client(blocks, image_width, image_height),
        }
        return _create_visual_pages_pptx([page])

    if source_kind == "pdf":
        pages = _extract_pdf_pages(content, updated_content)
        return _create_visual_pages_pptx(pages)

    return _convert_existing_pptx(content, updated_content)


def _normalize_image_bytes(content: bytes, content_type: str) -> Tuple[bytes, str, int, int]:
    with Image.open(io.BytesIO(content)) as opened:
        image = opened
        if getattr(image, "is_animated", False):
            image = next(ImageSequence.Iterator(image)).copy()
        else:
            image = image.copy()

    image = image.convert("RGB")
    output = io.BytesIO()
    image.save(output, format="PNG")
    return output.getvalue(), "image/png", image.width, image.height


def _bbox_from_bounds(x0: float, y0: float, x1: float, y1: float) -> List[List[int]]:
    return [
        [int(round(x0)), int(round(y0))],
        [int(round(x1)), int(round(y0))],
        [int(round(x1)), int(round(y1))],
        [int(round(x0)), int(round(y1))],
    ]


def _prepare_blocks_for_client(
    blocks: List[Dict[str, Any]],
    image_width: int,
    image_height: int,
) -> List[Dict[str, Any]]:
    prepared = []
    for index, block in enumerate(blocks):
        text = str(block.get("text", ""))
        prepared_block = {
            **block,
            "block_index": block.get("block_index", index),
            "text": text,
            "originalText": block.get("originalText", text),
            "image_width": block.get("image_width", image_width),
            "image_height": block.get("image_height", image_height),
        }
        prepared.append(prepared_block)
    return prepared


def _image_blocks_from_updated(updated_content: Any | None) -> List[Dict[str, Any]] | None:
    if updated_content is None:
        return None
    if isinstance(updated_content, dict) and isinstance(updated_content.get("blocks"), list):
        return updated_content["blocks"]
    if isinstance(updated_content, list):
        return updated_content
    return None


def _extract_pdf_pages(content: bytes, updated_content: Any | None = None) -> List[Dict[str, Any]]:
    document = fitz.open(stream=content, filetype="pdf")
    pages = []

    try:
        for page_index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(
                matrix=fitz.Matrix(PDF_RENDER_SCALE, PDF_RENDER_SCALE),
                alpha=False,
            )
            image_bytes = pixmap.tobytes("png")
            updated_blocks = _slide_blocks_from_updated(updated_content, page_index)

            if updated_blocks is None:
                blocks = _pdf_text_blocks(page, pixmap.width, pixmap.height)
                if not blocks:
                    ocr_result = extract_text_from_image(image_bytes)
                    if not ocr_result.get("success"):
                        blocks = []
                    else:
                        blocks = ocr_result.get("text_blocks", [])
            else:
                blocks = updated_blocks

            pages.append({
                "label": f"Page {page_index}",
                "image_bytes": image_bytes,
                "content_type": "image/png",
                "image_width": pixmap.width,
                "image_height": pixmap.height,
                "blocks": _prepare_blocks_for_client(blocks, pixmap.width, pixmap.height),
            })
    finally:
        document.close()

    return pages


def _pdf_text_blocks(page: fitz.Page, image_width: int, image_height: int) -> List[Dict[str, Any]]:
    text_dict = page.get_text("dict")
    blocks = []

    for text_block in text_dict.get("blocks", []):
        if text_block.get("type") != 0:
            continue

        for line in text_block.get("lines", []):
            for span in line.get("spans", []):
                text = str(span.get("text", "")).strip()
                if not text:
                    continue

                x0, y0, x1, y1 = span.get("bbox", (0, 0, 0, 0))
                pixel_bbox = _bbox_from_bounds(
                    x0 * PDF_RENDER_SCALE,
                    y0 * PDF_RENDER_SCALE,
                    x1 * PDF_RENDER_SCALE,
                    y1 * PDF_RENDER_SCALE,
                )
                flags = int(span.get("flags", 0) or 0)
                color = _color_from_pdf_int(int(span.get("color", 0) or 0))
                blocks.append({
                    "source": "pdf_text",
                    "text": text,
                    "originalText": text,
                    "bbox": pixel_bbox,
                    "line_bboxes": [pixel_bbox],
                    "alignment": _pdf_alignment((x0, x1), page.rect.width),
                    "font_family": _clean_font_name(str(span.get("font", "Arial"))),
                    "font_size": float(span.get("size", 12.0) or 12.0),
                    "font_color": color,
                    "bold": bool(flags & 16),
                    "italic": bool(flags & 2),
                    "image_width": image_width,
                    "image_height": image_height,
                })

    return blocks


def _color_from_pdf_int(value: int) -> Tuple[int, int, int]:
    return (value >> 16) & 255, (value >> 8) & 255, value & 255


def _clean_font_name(font_name: str) -> str:
    font_name = font_name.split("+")[-1]
    for suffix in ("-BoldItalic", "-BoldOblique", "-Bold", "-Italic", "-Oblique", "-Regular"):
        font_name = font_name.replace(suffix, "")
    return font_name.replace(",", " ").strip() or "Arial"


def _pdf_alignment(bounds: Tuple[float, float], page_width: float) -> str:
    x0, x1 = bounds
    center = (x0 + x1) / 2
    if abs(center - (page_width / 2)) <= page_width * 0.08:
        return "center"
    if center > page_width * 0.62:
        return "right"
    return "left"


def _page_to_client_slide(page: Dict[str, Any], page_number: int) -> Dict[str, Any]:
    full_text = "\n".join(block.get("text", "") for block in page.get("blocks", []) if block.get("text"))
    return {
        "slide_number": page_number,
        "text": full_text,
        "text_blocks": page.get("blocks", []),
        "full_text_with_label": f"Page {page_number}:\n{full_text}",
        "image_width": page.get("image_width"),
        "image_height": page.get("image_height"),
    }


def _slide_blocks_from_updated(updated_content: Any | None, slide_number: int) -> List[Dict[str, Any]] | None:
    if not isinstance(updated_content, dict):
        return None

    value = updated_content.get(str(slide_number), updated_content.get(slide_number))
    if isinstance(value, list):
        return value
    return None


def _shape_text_from_updated(updated_content: Any | None, slide_number: int, shape_index: int) -> str | None:
    blocks = _slide_blocks_from_updated(updated_content, slide_number)
    if blocks is None:
        return None

    for block in blocks:
        if block.get("source", "shape") != "shape":
            continue
        try:
            if int(block.get("shape_index")) == shape_index:
                return str(block.get("text", ""))
        except (TypeError, ValueError):
            continue

    return None


def _picture_blocks_for_shape(
    extracted_slides: List[Dict[str, Any]],
    updated_content: Any | None,
    slide_number: int,
    shape_index: int,
) -> List[Dict[str, Any]]:
    updated_blocks = _slide_blocks_from_updated(updated_content, slide_number)
    if updated_blocks is not None:
        return [
            block
            for block in updated_blocks
            if (block.get("source") == "picture" or block.get("bbox"))
            and _safe_int(block.get("shape_index")) == shape_index
        ]

    for slide in extracted_slides:
        if slide.get("slide_number") != slide_number:
            continue
        return [
            block
            for block in slide.get("text_blocks", [])
            if block.get("source") == "picture" and _safe_int(block.get("shape_index")) == shape_index
        ]
    return []


def _safe_int(value: Any) -> int | None:
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _create_visual_pages_pptx(pages: List[Dict[str, Any]]) -> bytes:
    if not pages:
        raise ValueError("No pages or images could be converted into slides.")

    presentation = Presentation()
    first_page = pages[0]
    slide_width, slide_height = _slide_size_for_aspect(
        int(first_page["image_width"]),
        int(first_page["image_height"]),
    )
    presentation.slide_width = slide_width
    presentation.slide_height = slide_height

    for page in pages:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _add_reconstructed_image_slide(slide, page, slide_width, slide_height)

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _slide_size_for_aspect(image_width: int, image_height: int) -> Tuple[Emu, Emu]:
    image_width = max(1, image_width)
    image_height = max(1, image_height)
    aspect = image_width / image_height

    if aspect >= 1:
        width_in = DEFAULT_LANDSCAPE_WIDTH_IN
        height_in = width_in / aspect
    else:
        height_in = DEFAULT_PORTRAIT_HEIGHT_IN
        width_in = height_in * aspect

    return Inches(width_in), Inches(height_in)


def _fit_image_to_slide(image_width: int, image_height: int, slide_width: int, slide_height: int) -> Tuple[int, int, int, int]:
    scale = min(slide_width / max(1, image_width), slide_height / max(1, image_height))
    width = int(image_width * scale)
    height = int(image_height * scale)
    left = int((slide_width - width) / 2)
    top = int((slide_height - height) / 2)
    return left, top, width, height


def _add_reconstructed_image_slide(slide, page: Dict[str, Any], slide_width: int, slide_height: int) -> None:
    image_bytes = page["image_bytes"]
    image_width = int(page["image_width"])
    image_height = int(page["image_height"])
    blocks = _prepare_blocks_for_client(page.get("blocks", []), image_width, image_height)
    clean_image = _clean_image_text(image_bytes, page.get("content_type", "image/png"), blocks)

    image_left, image_top, display_width, display_height = _fit_image_to_slide(
        image_width,
        image_height,
        slide_width,
        slide_height,
    )
    slide.shapes.add_picture(
        io.BytesIO(clean_image),
        image_left,
        image_top,
        width=display_width,
        height=display_height,
    )

    with Image.open(io.BytesIO(image_bytes)) as source_image:
        source_rgb = source_image.convert("RGB")
        for block in blocks:
            _add_textbox_for_image_block(
                slide,
                block,
                source_rgb,
                image_left,
                image_top,
                display_width,
                display_height,
                image_width,
                image_height,
            )


def _clean_image_text(
    image_bytes: bytes,
    content_type: str,
    blocks: List[Dict[str, Any]],
    normalize_output: bool = True,
) -> bytes:
    erase_blocks = []
    for block in blocks:
        original_text = str(block.get("originalText", block.get("text", "")))
        if not original_text.strip() or not block.get("bbox"):
            continue
        erase_blocks.append({
            **block,
            "text": "",
            "originalText": original_text,
        })

    if not erase_blocks:
        if not normalize_output:
            return image_bytes
        normalized, _, _, _ = _normalize_image_bytes(image_bytes, content_type)
        return normalized

    cleaned, _ = update_image_with_text(
        image_bytes,
        content_type if content_type else "image/png",
        {"blocks": erase_blocks},
    )
    if not normalize_output:
        return cleaned

    normalized, _, _, _ = _normalize_image_bytes(cleaned, content_type)
    return normalized


def _add_textbox_for_image_block(
    slide,
    block: Dict[str, Any],
    source_image: Image.Image,
    image_left: int,
    image_top: int,
    display_width: int,
    display_height: int,
    image_width: int,
    image_height: int,
) -> None:
    text = str(block.get("text", ""))
    if not text.strip() or not block.get("bbox"):
        return

    bbox = _bbox_bounds(block["bbox"])
    region = _expanded_text_region(block, bbox, (image_width, image_height))
    left, top, width, height = _map_pixel_box_to_slide(
        region,
        image_left,
        image_top,
        display_width,
        display_height,
        image_width,
        image_height,
    )
    _, sampled_text_color = _sample_block_colors(source_image, bbox)
    text_color = _block_color(block, sampled_text_color)
    line_count = _line_count_from_block(block, str(block.get("originalText", "")), text)
    font_size = _block_font_size_points(
        block,
        bbox,
        image_height,
        display_height,
        line_count,
    )
    alignment = _block_alignment(block, bbox, image_width)

    _add_ppt_textbox(
        slide,
        left,
        top,
        width,
        height,
        text,
        {
            "font_family": block.get("font_family") or "Arial",
            "font_size": font_size,
            "font_color": text_color,
            "alignment": alignment,
            "bold": bool(block.get("bold", False)),
            "italic": bool(block.get("italic", False)),
        },
    )


def _map_pixel_box_to_slide(
    box: Tuple[int, int, int, int],
    image_left: int,
    image_top: int,
    display_width: int,
    display_height: int,
    image_width: int,
    image_height: int,
) -> Tuple[int, int, int, int]:
    x0, y0, x1, y1 = box
    left = image_left + int((x0 / max(1, image_width)) * display_width)
    top = image_top + int((y0 / max(1, image_height)) * display_height)
    right = image_left + int((x1 / max(1, image_width)) * display_width)
    bottom = image_top + int((y1 / max(1, image_height)) * display_height)
    return left, top, max(1, right - left), max(1, bottom - top)


def _block_color(block: Dict[str, Any], fallback: Tuple[int, int, int]) -> Tuple[int, int, int]:
    color = block.get("font_color")
    if isinstance(color, (list, tuple)) and len(color) >= 3:
        return int(color[0]), int(color[1]), int(color[2])
    return fallback


def _block_font_size_points(
    block: Dict[str, Any],
    bbox: Tuple[int, int, int, int],
    image_height: int,
    display_height: int,
    line_count: int,
) -> float:
    if block.get("font_size") is not None:
        pdf_points = float(block.get("font_size") or 12)
        image_points = (display_height / EMU_PER_INCH) * 72
        source_points = max(1, image_height / PDF_RENDER_SCALE)
        if block.get("source") == "pdf_text":
            return _clamp(pdf_points * (image_points / source_points), 5, 96)

    _, y0, _, y1 = bbox
    pixel_line_height = max(1, (y1 - y0) / max(1, line_count))
    points = (pixel_line_height / max(1, image_height)) * (display_height / EMU_PER_INCH) * 72
    return _clamp(points * 0.9, 5, 96)


def _clamp(value: float, minimum: float, maximum: float) -> float:
    return max(minimum, min(maximum, value))


def _add_ppt_textbox(
    slide,
    left: int,
    top: int,
    width: int,
    height: int,
    text: str,
    style: Dict[str, Any],
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.margin_left = Emu(0)
    text_frame.margin_right = Emu(0)
    text_frame.margin_top = Emu(0)
    text_frame.margin_bottom = Emu(0)
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    lines = text.splitlines() or [text]
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = _ppt_alignment(style.get("alignment"))
        run = paragraph.add_run()
        run.text = line
        _apply_run_style(run, style)

    return textbox


def _ppt_alignment(alignment: Any):
    if str(alignment).lower() == "center":
        return PP_ALIGN.CENTER
    if str(alignment).lower() == "right":
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def _apply_run_style(run, style: Dict[str, Any]) -> None:
    font = run.font
    font.name = str(style.get("font_family") or "Arial")
    font.size = Pt(float(style.get("font_size") or 12))
    font.bold = bool(style.get("bold", False))
    font.italic = bool(style.get("italic", False))
    color = _block_color(style, (0, 0, 0))
    font.color.rgb = RGBColor(*color)


def _convert_existing_pptx(content: bytes, updated_content: Any | None) -> bytes:
    extracted_slides, presentation = process_pptx(content)

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                replacement_text = _shape_text_from_updated(updated_content, slide_number, shape_index)
                if replacement_text is not None:
                    _replace_shape_text(shape, replacement_text)
                continue

            if getattr(shape, "shape_type", None) != 13:
                continue

            blocks = _picture_blocks_for_shape(extracted_slides, updated_content, slide_number, shape_index)
            if not blocks:
                continue

            image_part = shape.part.related_part(shape._pic.blipFill.blip.rEmbed)
            image_content_type = shape.image.content_type
            original_image = image_part.blob
            clean_image = _clean_image_text(
                original_image,
                image_content_type,
                blocks,
                normalize_output=False,
            )
            image_part._blob = clean_image

            with Image.open(io.BytesIO(original_image)) as source_image:
                source_rgb = source_image.convert("RGB")
                image_width, image_height = source_rgb.size
                for block in _prepare_blocks_for_client(blocks, image_width, image_height):
                    _add_textbox_for_image_block(
                        slide,
                        block,
                        source_rgb,
                        int(shape.left),
                        int(shape.top),
                        int(shape.width),
                        int(shape.height),
                        int(block.get("image_width") or image_width),
                        int(block.get("image_height") or image_height),
                    )

    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()
