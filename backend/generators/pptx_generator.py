from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import base64
import io

COLORS = {
    "dark_blue": RGBColor(0x1F, 0x38, 0x64),
    "blue": RGBColor(0x44, 0x72, 0xC4),
    "light_blue": RGBColor(0xBD, 0xD7, 0xEE),
    "accent": RGBColor(0x70, 0xAD, 0x47),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "dark_gray": RGBColor(0x2E, 0x2E, 0x2E),
    "light_gray": RGBColor(0xF2, 0xF2, 0xF2),
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(blank_layout)


def _fill_shape(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    _fill_shape(shape, color)
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                 color: RGBColor = COLORS["dark_gray"], align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def _add_bullet_textbox(slide, left, top, width, height, bullets: list,
                        font_size=16, color: RGBColor = COLORS["dark_gray"]):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        run.text = f"• {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_before = Pt(4)

    return txBox


def _add_title_slide(prs: Presentation, info: dict):
    slide = _blank_slide(prs)

    # Full dark background
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["dark_blue"])

    # Accent bar on left
    _add_rect(slide, 0, 0, Inches(0.3), SLIDE_H, COLORS["blue"])

    # Bottom accent strip
    _add_rect(slide, 0, SLIDE_H - Inches(0.15), SLIDE_W, Inches(0.15), COLORS["accent"])

    # Title
    _add_textbox(
        slide,
        Inches(1), Inches(2.2), Inches(11), Inches(1.8),
        info.get("title", "Presentation Title"),
        font_size=44, bold=True, color=COLORS["white"], align=PP_ALIGN.LEFT,
    )

    # Subtitle
    subtitle = info.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide,
            Inches(1), Inches(4.2), Inches(11), Inches(1.2),
            subtitle,
            font_size=24, bold=False, color=COLORS["light_blue"], align=PP_ALIGN.LEFT,
        )

    _add_speaker_notes(slide, info.get("notes", ""))


def _add_content_slide(prs: Presentation, info: dict):
    slide = _blank_slide(prs)

    # Light background
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["light_gray"])

    # Header bar
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), COLORS["dark_blue"])

    # Accent line under header
    _add_rect(slide, 0, Inches(1.2), SLIDE_W, Inches(0.05), COLORS["blue"])

    # Title in header
    _add_textbox(
        slide,
        Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.9),
        info.get("title", ""),
        font_size=28, bold=True, color=COLORS["white"], align=PP_ALIGN.LEFT,
    )

    # Content card
    _add_rect(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.7), COLORS["white"])

    bullets = info.get("bullets", [])
    if bullets:
        _add_bullet_textbox(
            slide,
            Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.2),
            bullets, font_size=18,
        )

    _add_speaker_notes(slide, info.get("notes", ""))


def _add_two_column_slide(prs: Presentation, info: dict):
    slide = _blank_slide(prs)

    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["light_gray"])
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), COLORS["dark_blue"])
    _add_rect(slide, 0, Inches(1.2), SLIDE_W, Inches(0.05), COLORS["blue"])

    _add_textbox(
        slide,
        Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.9),
        info.get("title", ""),
        font_size=28, bold=True, color=COLORS["white"],
    )

    col_w = Inches(6.0)
    gap = Inches(0.4)
    col_h = Inches(5.5)
    top = Inches(1.45)

    # Left card
    _add_rect(slide, Inches(0.4), top, col_w, col_h, COLORS["white"])
    _add_rect(slide, Inches(0.4), top, col_w, Inches(0.45), COLORS["blue"])
    left_title = info.get("left_title", "")
    if left_title:
        _add_textbox(slide, Inches(0.55), top + Pt(4), col_w - Inches(0.2), Inches(0.4),
                     left_title, font_size=14, bold=True, color=COLORS["white"])
    left_bullets = info.get("left_bullets", [])
    if left_bullets:
        _add_bullet_textbox(slide, Inches(0.55), top + Inches(0.55), col_w - Inches(0.2),
                            col_h - Inches(0.65), left_bullets, font_size=16)

    # Right card
    right_left = Inches(0.4) + col_w + gap
    _add_rect(slide, right_left, top, col_w, col_h, COLORS["white"])
    _add_rect(slide, right_left, top, col_w, Inches(0.45), COLORS["dark_blue"])
    right_title = info.get("right_title", "")
    if right_title:
        _add_textbox(slide, right_left + Inches(0.15), top + Pt(4),
                     col_w - Inches(0.2), Inches(0.4),
                     right_title, font_size=14, bold=True, color=COLORS["white"])
    right_bullets = info.get("right_bullets", [])
    if right_bullets:
        _add_bullet_textbox(slide, right_left + Inches(0.15), top + Inches(0.55),
                            col_w - Inches(0.2), col_h - Inches(0.65), right_bullets, font_size=16)

    _add_speaker_notes(slide, info.get("notes", ""))


def _add_section_slide(prs: Presentation, info: dict):
    slide = _blank_slide(prs)

    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["blue"])
    _add_rect(slide, 0, SLIDE_H - Inches(0.12), SLIDE_W, Inches(0.12), COLORS["accent"])

    _add_textbox(
        slide,
        Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
        info.get("title", ""),
        font_size=40, bold=True, color=COLORS["white"], align=PP_ALIGN.CENTER,
    )

    subtitle = info.get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide,
            Inches(1.5), Inches(4.2), Inches(10), Inches(1.0),
            subtitle,
            font_size=20, bold=False, color=COLORS["light_blue"], align=PP_ALIGN.CENTER,
        )

    _add_speaker_notes(slide, info.get("notes", ""))


def _add_speaker_notes(slide, notes_text: str):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def _fit_box(img_w, img_h, box_left, box_top, box_w, box_h):
    scale = min(box_w / img_w, box_h / img_h)
    width = int(img_w * scale)
    height = int(img_h * scale)
    left = int(box_left + (box_w - width) / 2)
    top = int(box_top + (box_h - height) / 2)
    return left, top, width, height


def _add_source_visual_slide(prs: Presentation, source: dict, index: int, total: int):
    slide = _blank_slide(prs)

    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, COLORS["light_gray"])
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), COLORS["dark_blue"])
    _add_rect(slide, 0, Inches(0.85), SLIDE_W, Inches(0.05), COLORS["blue"])

    title = source.get("label") or source.get("filename") or f"Source Content {index}"
    _add_textbox(
        slide,
        Inches(0.4), Inches(0.12), Inches(9.5), Inches(0.55),
        title,
        font_size=20, bold=True, color=COLORS["white"], align=PP_ALIGN.LEFT,
    )
    _add_textbox(
        slide,
        Inches(10.4), Inches(0.17), Inches(2.5), Inches(0.45),
        f"{index} of {total}",
        font_size=12, bold=False, color=COLORS["light_blue"], align=PP_ALIGN.RIGHT,
    )

    image_bytes = base64.b64decode(source["base64"])
    image_stream = io.BytesIO(image_bytes)
    with Image.open(io.BytesIO(image_bytes)) as img:
        img_w, img_h = img.size

    box_left = Inches(0.55)
    box_top = Inches(1.12)
    box_w = SLIDE_W - Inches(1.1)
    box_h = SLIDE_H - Inches(1.55)
    left, top, width, height = _fit_box(img_w, img_h, box_left, box_top, box_w, box_h)

    slide.shapes.add_picture(image_stream, left, top, width=width, height=height)
    _add_speaker_notes(slide, f"Original uploaded content included in the presentation: {title}")


def _add_source_visuals(prs: Presentation, source_visuals: list):
    if not source_visuals:
        return

    _add_section_slide(prs, {
        "title": "Source Content",
        "subtitle": "Original uploaded files and PDF page previews",
        "notes": "These slides preserve the uploaded source content inside the generated presentation.",
    })
    total = len(source_visuals)
    for index, source in enumerate(source_visuals, start=1):
        _add_source_visual_slide(prs, source, index, total)


def create_presentation(slide_data: dict, source_visuals: list | None = None) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for slide_info in slide_data.get("slides", []):
        slide_type = slide_info.get("type", "content")
        if slide_type == "title":
            _add_title_slide(prs, slide_info)
        elif slide_type == "two_column":
            _add_two_column_slide(prs, slide_info)
        elif slide_type == "section":
            _add_section_slide(prs, slide_info)
        else:
            _add_content_slide(prs, slide_info)

    _add_source_visuals(prs, source_visuals or [])

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()
